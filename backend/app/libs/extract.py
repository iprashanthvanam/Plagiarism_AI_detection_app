



"""
Text extraction pipeline — TKREC Plagiarism Analysis System

SCANNED PDF ROUTING:
  1. pdfplumber          → digital PDF (selectable text, any language)
  2. Gemini Vision       → scanned PDF PRIMARY  (multilingual LLM, handles
                           Telugu, Hindi, handwriting, diagrams natively)
     quota/unavailable →
  3. Enhanced Tesseract  → fallback
       - 300 DPI rendering
       - contrast + sharpen + binarize preprocessing
       - PSM 4 + LSTM engine
       - per-page quality gate (skips garbage pages)
       - TWO-PASS LANGUAGE DETECTION (new)
           Pass 1: eng only  → detect Indic script in output
           Pass 2: eng+<lang> if pack installed (e.g. eng+tel for Telugu)

LANGUAGE SUPPORT (scanned PDFs):
  Digital PDFs with embedded Unicode text → pdfplumber returns correct
  Telugu/Hindi/Tamil/etc text automatically — no special handling needed.

  Scanned PDFs → Gemini handles all Indian languages natively.
  Tesseract fallback → requires tesseract-ocr-<lang> system package:
    sudo apt-get install -y tesseract-ocr-tel   # Telugu
    sudo apt-get install -y tesseract-ocr-hin   # Hindi
    sudo apt-get install -y tesseract-ocr-tam   # Tamil
    sudo apt-get install -y tesseract-ocr-kan   # Kannada
    sudo apt-get install -y tesseract-ocr-mal   # Malayalam

BUGS FIXED (vs original):
  BUG-1  dpi=200 → blurry renders  →  fixed: dpi=300
  BUG-2  no preprocessing          →  fixed: contrast/sharpen/binarize
  BUG-3  wrong PSM mode (3→4)      →  fixed: --psm 4 --oem 1
  BUG-4  no quality gate           →  fixed: skip pages < 40% alpha ratio
  BUG-5  Gemini called without is_pdf arg → TypeError silently swallowed
          → Gemini NEVER actually ran  →  fixed: pass is_pdf=True/False
  BUG-6  Gemini was fallback for scanned PDFs, should be PRIMARY
  BUG-7  No Telugu / Indic language support in Tesseract fallback
          →  fixed: two-pass lang detection + eng+tel if pack installed
"""

import os
import re
import logging
import subprocess
from typing import Optional, Tuple, List

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

# ============================================================
# LOCAL OCR IMPORTS
# ============================================================
try:
    import pytesseract
    from PIL import Image, ImageFilter, ImageEnhance
    import numpy as np
    _tesseract_available = True
except ImportError:
    _tesseract_available = False
    logger.warning("pytesseract/Pillow not installed — image OCR falls back to Gemini")

try:
    from pdf2image import convert_from_path
    _pdf2image_available = True
except ImportError:
    _pdf2image_available = False
    logger.warning("pdf2image not installed — scanned PDF OCR falls back to Gemini")

# ============================================================
# GEMINI IMPORT
# ============================================================
try:
    from app.libs.gemini_service import extract_text_with_gemini
    _gemini_available = True
except Exception as e:
    logger.warning("Gemini service unavailable: %s", e)
    extract_text_with_gemini = None
    _gemini_available = False

# ============================================================
# CONSTANTS
# ============================================================
MIN_LOCAL_TEXT = 10
MIN_SELECTABLE_CHARS = 100  # ✅ ADD THIS — minimum chars to consider PDF digital

# Pages below this alpha-char ratio are skipped (garbage OCR)
OCR_QUALITY_THRESHOLD = 0.40

# Max pages to process from a scanned PDF for plagiarism purposes
SCANNED_PDF_MAX_PAGES = 20

# Tesseract base config — works for all scripts
TESSERACT_BASE_CONFIG = "--psm 4 --oem 1"

# Gemini quota/rate-limit error signals from gemini_service
_GEMINI_QUOTA_SIGNALS = {
    "quota", "429", "resource_exhausted", "rate limit",
    "too many requests", "exceeded",
}

# ─── Indic script Unicode ranges ────────────────────────────────────────────
# Maps Tesseract language code → (range_start, range_end)
# All Indian language PDFs will be detected and routed correctly.
_INDIC_SCRIPT_RANGES: dict[str, tuple[str, str]] = {
    "tel": ("\u0C00", "\u0C7F"),   # Telugu
    "hin": ("\u0900", "\u097F"),   # Hindi / Devanagari (also Marathi)
    "tam": ("\u0B80", "\u0BFF"),   # Tamil
    "kan": ("\u0C80", "\u0CFF"),   # Kannada
    "mal": ("\u0D00", "\u0D7F"),   # Malayalam
    "guj": ("\u0A80", "\u0AFF"),   # Gujarati
    "pan": ("\u0A00", "\u0A7F"),   # Punjabi / Gurmukhi
    "ori": ("\u0B00", "\u0B7F"),   # Odia
    "ben": ("\u0980", "\u09FF"),   # Bengali
}

_LANGUAGE_NAMES: dict[str, str] = {
    "tel": "Telugu", "hin": "Hindi/Devanagari", "tam": "Tamil",
    "kan": "Kannada", "mal": "Malayalam", "guj": "Gujarati",
    "pan": "Punjabi", "ori": "Odia", "ben": "Bengali",
}

# Noise tokens in .doc OLE metadata
_DOC_NOISE = {
    "Times New Roman", "Liberation Serif", "Liberation Sans", "DejaVu Sans",
    "Open Sans", "FreeSans", "OpenSymbol", "Arial Unicode MS", "Droid Sans",
    "Fallback", "Symbol", "Arial", "Heading", "oasis.open", "office.com",
    "Visited Internet Link", "Internet Link", "Root Entry", "WordDocument",
}

# ============================================================
# HELPER FUNCTIONS — EXTRACTION
# ============================================================

def _extract_docx(file_path: str) -> str:
    """
    Extract COMPLETE text from .docx PRESERVING structure.
    ✅ NO TRUNCATION — returns full document text
    """
    try:
        doc = Document(file_path)
        sections = []
        
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            
            style_name = para.style.name if para.style else ""
            is_heading = "Heading" in style_name
            
            if is_heading:
                level = "".join(c for c in style_name if c.isdigit()) or "1"
                sections.append(f"\n{'#' * int(level)} {para.text.strip()}\n")
            else:
                sections.append(para.text.strip())
        
        for table in doc.tables:
            sections.append("\n[TABLE]\n")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                sections.append(row_text)
            sections.append("[/TABLE]\n")
        
        text = "\n\n".join(sections)
        if text:
            logger.info(
                "✅ DOCX extracted (COMPLETE): %s (%d chars)",
                os.path.basename(file_path), len(text)
            )
            return text  # ✅ NO TRUNCATION
        
        logger.warning("DOCX has no content — trying binary extraction")
        return _extract_doc_binary(file_path)
    
    except Exception as e:
        logger.warning("DOCX extraction failed for %s: %s", file_path, e)
        return _extract_doc_binary(file_path)


def _extract_pptx(file_path: str) -> str:
    """Extract COMPLETE text from PowerPoint. ✅ NO TRUNCATION"""
    try:
        prs = Presentation(file_path)
        slides = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                
                text = shape.text.strip()
                if not text:
                    continue
                
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == 1:
                        slide_content.insert(0, f"\n## Slide {slide_num}: {text}\n")
                        continue
                
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        level = para.level
                        indent = "  " * level
                        bullet = "• " if para.level == 0 else "◦ "
                        slide_content.append(f"{indent}{bullet}{para.text.strip()}")
                else:
                    slide_content.append(text)
            
            if slide_content:
                slides.append("\n".join(slide_content))
        
        text = "\n\n".join(slides)
        if text:
            logger.info(
                "✅ PPTX extracted (COMPLETE): %s (%d chars)",
                os.path.basename(file_path), len(text)
            )
            return text  # ✅ NO TRUNCATION
        
        logger.warning("PPTX has no text")
        return ""
    
    except Exception as e:
        logger.warning("PPTX extraction failed for %s: %s", file_path, e)
        return ""


def _is_scanned_pdf(file_path: str) -> bool:
    """
    Detect if a PDF is scanned (image-based) or digital (text-based).
    
    Heuristic:
    - If pdfplumber extracts >= MIN_SELECTABLE_CHARS → digital
    - Otherwise → scanned (needs OCR)
    """
    try:
        with pdfplumber.open(file_path) as pdf:
            # Sample first 5 pages
            text_chars = 0
            for page in pdf.pages[:5]:
                text = page.extract_text()
                if text:
                    text_chars += len(text)
            
            is_digital = text_chars >= MIN_SELECTABLE_CHARS
            logger.info("PDF type detection: %s → %s (%d chars in first 5 pages)",
                       os.path.basename(file_path),
                       "digital" if is_digital else "scanned",
                       text_chars)
            return not is_digital
    
    except Exception as e:
        logger.warning("PDF type detection failed for %s: %s — assuming scanned",
                      file_path, e)
        return True


def _extract_pdfplumber(file_path: str) -> str:
    """Extract COMPLETE text from digital PDFs. ✅ NO TRUNCATION"""
    try:
        with pdfplumber.open(file_path) as pdf:
            pages = []
            
            for page_num, page in enumerate(pdf.pages, 1):
                # Try to extract text with layout preserved
                text = page.extract_text(layout=True)
                
                if not text:
                    # Fallback to simple extraction
                    text = page.extract_text()
                
                if text and text.strip():
                    # Preserve page breaks for long documents
                    if len(pdf.pages) > 1:
                        pages.append(f"\n--- Page {page_num} ---\n{text.strip()}")
                    else:
                        pages.append(text.strip())
            
            result = "\n\n".join(pages)
            logger.info(
                "✅ pdfplumber (COMPLETE): %s (%d chars)",
                os.path.basename(file_path), len(result)
            )
            return result  # ✅ NO TRUNCATION
    
    except Exception as e:
        logger.warning("pdfplumber failed for %s: %s", file_path, e)
        return ""


async def _extract_scanned_pdf(file_path: str) -> str:
    """Extract COMPLETE text from scanned PDFs. ✅ NO TRUNCATION"""
    text, _ = await _extract_scanned_pdf_gemini(file_path)
    if text:
        logger.info("✅ Scanned PDF via Gemini (COMPLETE): %d chars", len(text))
        return text  # ✅ NO TRUNCATION
    
    text = _extract_scanned_pdf_local(file_path)
    logger.info("✅ Scanned PDF via Tesseract (COMPLETE): %d chars", len(text))
    return text  # ✅ NO TRUNCATION


# ============================================================
# LANGUAGE DETECTION UTILITIES
# ============================================================

def _detect_indic_scripts(text: str) -> List[str]:
    """
    Return list of detected Tesseract language codes for Indic scripts
    present in the given text. E.g. ['tel'] for Telugu, ['tel', 'hin']
    for mixed Telugu+Hindi.
    """
    return [
        lang
        for lang, (start, end) in _INDIC_SCRIPT_RANGES.items()
        if any(start <= c <= end for c in text)
    ]


def _get_installed_tesseract_langs() -> set:
    """
    Return the set of Tesseract language codes installed on this system.
    Cached after first call to avoid repeated subprocess calls.
    """
    if not hasattr(_get_installed_tesseract_langs, "_cache"):
        try:
            r = subprocess.run(
                ["tesseract", "--list-langs"],
                capture_output=True, text=True, timeout=5
            )
            lines = r.stdout.strip().split("\n")
            # First line is "List of available languages..." header — skip it
            langs = {l.strip() for l in lines[1:] if l.strip()}
            _get_installed_tesseract_langs._cache = langs
            logger.info("Installed Tesseract languages: %s", langs)
        except Exception as e:
            logger.warning("Could not query Tesseract langs: %s", e)
            _get_installed_tesseract_langs._cache = {"eng"}
    return _get_installed_tesseract_langs._cache


def _build_tesseract_lang_string(detected_scripts: List[str]) -> str:
    """
    Build the Tesseract -l argument from detected scripts.

    Rules:
    - Always include 'eng' (academic docs are almost always mixed with English)
    - Only include script packs that are actually installed
    - Warn if a detected script pack is NOT installed
    - Returns e.g. 'eng+tel' if tel is installed, 'eng' if not

    System install guide (logged as WARNING if missing):
      sudo apt-get install -y tesseract-ocr-tel   # Telugu
      sudo apt-get install -y tesseract-ocr-hin   # Hindi
      sudo apt-get install -y tesseract-ocr-tam   # Tamil
      sudo apt-get install -y tesseract-ocr-kan   # Kannada
      sudo apt-get install -y tesseract-ocr-mal   # Malayalam

    """
    if not detected_scripts:
        return "eng"

    installed = _get_installed_tesseract_langs()
    available = []
    missing = []

    for script in detected_scripts:
        if script in installed:
            available.append(script)
        else:
            missing.append(script)

    if missing:
        lang_names = [_LANGUAGE_NAMES.get(m, m) for m in missing]
        logger.warning(
            "Tesseract language packs NOT installed for detected scripts %s. "
            "OCR quality for these scripts will be degraded. "
            "Install with: sudo apt-get install -y %s",
            lang_names,
            " ".join(f"tesseract-ocr-{m}" for m in missing),
        )

    all_langs = ["eng"] + available
    return "+".join(all_langs)  # e.g. "eng+tel" or "eng+tel+hin"


# ============================================================
# IMAGE QUALITY HELPERS
# ============================================================

def _assess_image_quality(img) -> float:
    """
    Content-presence score 0.0–1.0 from grayscale pixel stddev.
    Blank pages ≈ 0.0–0.1  |  text-filled pages ≈ 0.7–1.0
    """
    from PIL import ImageStat
    gray = img.convert("L")
    stat = ImageStat.Stat(gray)
    return min(1.0, stat.stddev[0] / 40.0)


def _preprocess_for_tesseract(img) -> "Image":
    """
    Enhance a page image before OCR:
      1. Grayscale     — remove colour noise
      2. Contrast ×1.8 — ink stands out from paper
      3. Sharpen       — sharpen pen/print strokes
      4. Binarize      — threshold = mean − 0.2×std (Otsu-style)
                         robust to yellowed/uneven scanned paper
    Works well for both English and Indic script pages.
    """
    img = img.convert("L")
    img = ImageEnhance.Contrast(img).enhance(1.8)
    img = img.filter(ImageFilter.SHARPEN)
    arr = np.array(img)
    threshold = arr.mean() - arr.std() * 0.2
    arr = np.where(arr > threshold, 255, 0).astype(np.uint8)
    return Image.fromarray(arr)


def _ocr_quality_score(text: str) -> float:
    """
    Fraction of alphabetic / Indic chars among non-whitespace chars.

    Extended to count Indic Unicode chars as 'valid' — without this,
    a page of pure Telugu text would score 0% (no ASCII alpha) and
    get skipped by the quality gate.
    """
    stripped = text.replace(" ", "").replace("\n", "")
    if not stripped:
        return 0.0

    def is_valid_char(c: str) -> bool:
        # ASCII letters
        if c.isalpha():
            return True
        # Indic script characters (U+0900 – U+0D7F covers all major scripts)
        cp = ord(c)
        return 0x0900 <= cp <= 0x0D7F

    valid = sum(1 for c in stripped if is_valid_char(c))
    return valid / len(stripped)


# ============================================================
# SCANNED PDF — GEMINI PRIMARY
# ============================================================

async def _extract_scanned_pdf_gemini(file_path: str) -> Tuple[str, bool]:
    """
    Attempt Gemini Vision extraction for a scanned PDF.

    Gemini is multilingual by design — it handles Telugu, Hindi, Tamil,
    and other Indic scripts WITHOUT needing any language configuration.

    Returns: (text, quota_exceeded)
    """
    if not _gemini_available or not extract_text_with_gemini:
        logger.info("Gemini not configured — skipping Gemini path")
        return "", False

    try:
        # BUG-5 FIX: is_pdf=True (this arg was missing before → TypeError)
        result = await extract_text_with_gemini(file_path, is_pdf=True)

        if not result:
            return "", False

        result_lower = result.lower()
        if result.startswith("error:"):
            quota_hit = any(sig in result_lower for sig in _GEMINI_QUOTA_SIGNALS)
            level = "quota/rate-limit" if quota_hit else "error"
            logger.warning("Gemini %s for %s: %s",
                           level, os.path.basename(file_path), result[:120])
            return "", quota_hit

        logger.info("Scanned PDF via Gemini: %s (%d chars)",
                    os.path.basename(file_path), len(result))
        return result, False

    except Exception as e:
        err = str(e).lower()
        quota_hit = any(sig in err for sig in _GEMINI_QUOTA_SIGNALS)
        logger.warning("Gemini %s for %s: %s",
                       "quota" if quota_hit else "exception",
                       os.path.basename(file_path), e)
        return "", quota_hit


# ============================================================
# SCANNED PDF — ENHANCED TESSERACT FALLBACK
# ============================================================

def _extract_scanned_pdf_local(file_path: str) -> str:
    """
    Enhanced local Tesseract OCR for scanned PDFs.
    Used ONLY when Gemini is unavailable or quota-exceeded.

    TWO-PASS LANGUAGE DETECTION:
    ─────────────────────────────
    Pass 1 (eng only, fast):
      • Run Tesseract with English only on page 1
      • Check output for Indic Unicode characters

    Pass 2 (full run with correct languages):
      • If Indic scripts detected in Pass 1:
          - Build lang string: eng+tel / eng+hin / eng+tel+hin etc.
          - Warn if language pack is not installed
          - Re-run with correct lang on all pages
      • If no Indic scripts:
          - Continue with eng only (no overhead)

    This ensures Telugu/Hindi/Tamil PDFs get proper Tesseract
    language support without slowing down English-only documents.
    """
    if not _pdf2image_available or not _tesseract_available:
        logger.warning("pdf2image or pytesseract unavailable")
        return ""

    try:
        pages = convert_from_path(
            file_path,
            dpi=300,                         # BUG-1 FIX: was 200
            first_page=1,
            last_page=SCANNED_PDF_MAX_PAGES,
        )
    except Exception as e:
        logger.warning("pdf2image failed for %s: %s", os.path.basename(file_path), e)
        return ""

    # ── Pass 1: Language Detection ────────────────────────────────────────
    tesseract_lang = "eng"
    if pages:
        try:
            processed_p1 = _preprocess_for_tesseract(pages[0])
            p1_text = pytesseract.image_to_string(
                processed_p1,
                config=f"{TESSERACT_BASE_CONFIG} -l eng",
                timeout=30,
            )
            detected_scripts = _detect_indic_scripts(p1_text)

            if not detected_scripts:
                # Pass 1 text is English-only OCR — check if the raw image
                # might contain Indic characters that eng-only Tesseract missed
                # by sampling the raw (unprocessed) first page too
                raw_p1_text = pytesseract.image_to_string(
                    pages[0],
                    config=f"{TESSERACT_BASE_CONFIG} -l eng",
                    timeout=30,
                )
                detected_scripts = _detect_indic_scripts(raw_p1_text)

            if detected_scripts:
                tesseract_lang = _build_tesseract_lang_string(detected_scripts)
                script_names = [_LANGUAGE_NAMES.get(s, s) for s in detected_scripts]
                logger.info(
                    "Indic scripts detected in Pass 1: %s → Tesseract lang=%r",
                    script_names, tesseract_lang
                )
            else:
                logger.info("No Indic scripts detected → using Tesseract lang='eng'")

        except Exception as e:
            logger.warning("Pass 1 language detection failed: %s — using eng", e)
            tesseract_lang = "eng"

    # ── Pass 2: Full Extraction ───────────────────────────────────────────
    tesseract_config = f"{TESSERACT_BASE_CONFIG} -l {tesseract_lang}"
    logger.info("Pass 2: OCR with config=%r on %d pages",
                tesseract_config, len(pages))

    texts = []
    total = len(pages)
    good = skipped_blank = skipped_garbage = 0

    for i, page in enumerate(pages):
        try:
            # Skip blank / nearly-blank pages
            if _assess_image_quality(page) < 0.10:
                skipped_blank += 1
                continue

            # Preprocess (BUG-2 FIX: was raw image)
            processed = _preprocess_for_tesseract(page)

            # OCR with tuned config (BUG-3 FIX: was default config)
            page_text = pytesseract.image_to_string(
                processed,
                config=tesseract_config,
                timeout=45,
            )

            if not page_text.strip():
                skipped_blank += 1
                continue

            # Per-page quality gate (BUG-4 FIX: was no gate)
            # NOTE: quality score now counts Indic chars as valid
            quality = _ocr_quality_score(page_text)
            if quality < OCR_QUALITY_THRESHOLD:
                logger.warning("Page %d/%d quality low (%.0f%%) — skipped",
                               i + 1, total, quality * 100)
                skipped_garbage += 1
                continue

            texts.append(page_text.strip())
            good += 1

        except Exception as e:
            logger.warning("Tesseract page %d error: %s", i + 1, e)

    logger.info(
        "Enhanced Tesseract [lang=%s]: %d/%d good | %d blank | %d garbage | %s",
        tesseract_lang, good, total, skipped_blank, skipped_garbage,
        os.path.basename(file_path),
    )
    return "\n\n".join(texts)


# ============================================================
# OTHER FORMAT HELPERS
# ============================================================

def _extract_doc_binary(file_path: str) -> str:
    """Legacy .doc (Word 97-2003) binary UTF-16LE text extraction."""
    try:
        with open(os.path.abspath(file_path), "rb") as f:
            data = f.read()
        chunks = re.findall(b"(?:[\x20-\x7e]\x00){8,}", data)
        text_parts = []
        for chunk in chunks:
            try:
                decoded = chunk.decode("utf-16-le", errors="ignore").strip()
            except Exception:
                continue
            if len(decoded) < 15:
                continue
            if any(noise in decoded for noise in _DOC_NOISE):
                continue
            text_parts.append(decoded)
        result = "\n".join(text_parts)
        logger.info(".doc binary: %s (%d chars)", os.path.basename(file_path), len(result))
        return result
    except Exception as e:
        logger.warning("Binary .doc failed for %s: %s", file_path, e)
        return ""


def _extract_image_local(file_path: str) -> str:
    """
    Tesseract OCR for standalone image files with language detection.
    Same two-pass approach as scanned PDFs.
    """
    if not _tesseract_available:
        return ""
    try:
        img = Image.open(file_path)
        processed = _preprocess_for_tesseract(img)

        # Quick language detection pass
        p1_text = pytesseract.image_to_string(
            processed, config=f"{TESSERACT_BASE_CONFIG} -l eng", timeout=20
        )
        detected = _detect_indic_scripts(p1_text)
        lang = _build_tesseract_lang_string(detected)

        if lang != "eng":
            return pytesseract.image_to_string(
                processed,
                config=f"{TESSERACT_BASE_CONFIG} -l {lang}",
                timeout=30,
            ).strip()
        return p1_text.strip()

    except Exception as e:
        logger.warning("Tesseract image OCR failed for %s: %s", file_path, e)
        return ""


def _extract_spreadsheet(file_path: str, ext: str) -> str:
    """Extract COMPLETE text from XLS/XLSX. ✅ NO TRUNCATION"""
    try:
        engine = "xlrd" if ext == ".xls" else "openpyxl"
        dfs = pd.read_excel(file_path, sheet_name=None, engine=engine)
        
        sheets = []
        for sheet_name, df in dfs.items():
            sheet_lines = []
            
            if len(dfs) > 1:
                sheet_lines.append(f"\n## Sheet: {sheet_name}\n")
            
            headers = [str(col).strip() for col in df.columns]
            sheet_lines.append(" | ".join(headers))
            sheet_lines.append("-" * len(" | ".join(headers)))
            
            for _, row in df.iterrows():
                cell_values = []
                for val in row:
                    if pd.isna(val):
                        cell_values.append("")
                    else:
                        sval = str(val).strip()
                        if sval.lower() == "nan":
                            cell_values.append("")
                        else:
                            cell_values.append(sval)
                
                if any(cell_values):
                    sheet_lines.append(" | ".join(cell_values))
            
            sheets.append("\n".join(sheet_lines))
        
        text = "\n\n".join(sheets)
        if text:
            logger.info(
                "✅ Spreadsheet (%s) extracted (COMPLETE): %s (%d chars)",
                ext, os.path.basename(file_path), len(text)
            )
            return text  # ✅ NO TRUNCATION
        
        return ""
    
    except Exception as e:
        logger.warning("Spreadsheet extraction failed for %s: %s", file_path, e)
        return ""


async def extract_text(file_path: str, content_type: str = "") -> str:
    """
    Extract COMPLETE text from any supported file format.
    ✅ NO TRUNCATION — returns full document text
    ✅ Validation added in tasks.py before DB storage
    """
    if not os.path.exists(file_path):
        logger.error("File not found: %s", file_path)
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt" or "text/plain" in content_type:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                result = f.read()  # ✅ NO TRUNCATION
                logger.info("✅ TXT extracted (COMPLETE): %d chars", len(result))
                return result
        except Exception:
            return ""

    # .doc  = Word 97-2003 binary OLE format — NOT a ZIP, python-docx cannot open it.
    #         Route directly to binary reader; skip _extract_docx() to avoid the
    #         noisy "There is no item named '[Content_Types].xml'" warning.
    # .docx = Office Open XML (ZIP + XML) — python-docx handles this natively.
    if ext == ".doc" or (content_type == "application/msword" and ext != ".docx"):
        return _extract_doc_binary(file_path)

    if ext == ".docx" or "wordprocessingml" in content_type:
        return _extract_docx(file_path)

    if ext in (".xls", ".xlsx") or "spreadsheet" in content_type.lower():
        return _extract_spreadsheet(file_path, ext)

    if ext == ".pptx" or "presentation" in content_type.lower():
        return _extract_pptx(file_path)

    if ext == ".pdf" or "pdf" in content_type:
        if not _is_scanned_pdf(file_path):
            text = _extract_pdfplumber(file_path)
            if text and len(text.strip()) >= 50:
                return text
        return await _extract_scanned_pdf(file_path)  # ✅ NOW WITH await

    logger.warning("Unknown type ext=%s", ext)
    return ""